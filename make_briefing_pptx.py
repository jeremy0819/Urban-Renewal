#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
都市更新／危老重建 住戶說明會簡報範本 → PowerPoint 產生器
風格參考 jieceng-web（暖白 × 炭黑 × 翡翠綠，Inter／微軟正黑，16:9）。
通用範本：所有個案數字以〔…〕佔位，無任何真實案件資料。

用法：  python3 make_briefing_pptx.py
輸出：  都更說明會簡報範本.pptx
相依：  pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches as IN, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 品牌色票 ----
WARM   = RGBColor.from_string("FAFAF9")  # 暖白底
PAPER  = RGBColor.from_string("F5F3EE")  # 暖白（卡片底）
CHAR   = RGBColor.from_string("1C1A17")  # 炭黑（標題）
BODY   = RGBColor.from_string("44403C")  # 內文
STONE  = RGBColor.from_string("A8A29E")  # 淺灰／佔位
LINE   = RGBColor.from_string("E7E5E4")  # 細線
WHITE  = RGBColor.from_string("FFFFFF")
EMER   = RGBColor.from_string("064E3B")  # 翡翠綠（品牌主色）
EMER_L = RGBColor.from_string("10B981")  # 亮翡翠（強調）
MINT   = RGBColor.from_string("ECFDF5")  # 極淺薄荷
MIST   = RGBColor.from_string("D8E6DF")  # 深底上的淺字

SANS, EA = "Inter", "Microsoft JhengHei"   # 拉丁 / 中日韓
SERIF    = "Noto Serif TC"

EMW, EMH = IN(13.333), IN(7.5)


def _ea(run, ea=EA):
    """設定 run 的東亞字型（python-pptx 預設只設拉丁）。"""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", ea)


def text(slide, l, t, w, h, runs, size=18, color=BODY, font=SANS, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.04, ls=None):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, {})]
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = sp
    for txt, o in runs:
        r = p.add_run(); r.text = txt
        r.font.size = Pt(o.get("size", size)); r.font.bold = o.get("bold", bold)
        r.font.name = o.get("font", font); r.font.color.rgb = o.get("color", color)
        if ls is not None or "ls" in o:
            r.font._rPr.set("spc", str(int((o.get("ls", ls) or 0) * 100)))
        _ea(r, o.get("ea", EA))
    return tb


def rect(slide, l, t, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp


def blank(prs, bg=WARM):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, EMW, EMH, fill=bg)
    return s


def frame(slide, color):
    """細框線（呼應 HTML 版的 hairline frame）。"""
    f = rect(slide, IN(0.42), IN(0.32), EMW - IN(0.84), EMH - IN(0.64))
    f.fill.background(); f.line.color.rgb = color; f.line.width = Pt(0.75)


def header(slide, kicker, title, num=None):
    text(slide, IN(0.9), IN(0.78), IN(11), IN(0.4), kicker.upper(),
         size=12, color=EMER_L, font=SANS, bold=True, ls=2.6)
    rect(slide, IN(0.9), IN(1.28), IN(0.7), Pt(3), fill=EMER)
    text(slide, IN(0.9), IN(1.46), IN(11.5), IN(1.0), title,
         size=33, color=CHAR, font=SERIF, bold=True)
    if num:
        text(slide, EMW - IN(2.2), IN(0.62), IN(1.6), IN(1.2), num,
             size=40, color=RGBColor.from_string("E4E0D6"), font=SERIF, bold=True,
             align=PP_ALIGN.RIGHT)


def footer(slide, label):
    text(slide, IN(0.9), EMH - IN(0.62), IN(8), IN(0.3), label,
         size=10.5, color=STONE, font=SANS, ls=0.8)


def cards(slide, items, top=IN(2.7), cols=3, gap=IN(0.32),
          left=IN(0.9), right=IN(0.9), rh=IN(1.78)):
    avail = EMW - left - right
    cw = (avail - gap * (cols - 1)) / cols
    rows = (len(items) + cols - 1) // cols
    for i, (emoji, title, bodytx) in enumerate(items):
        r, c = divmod(i, cols)
        x = left + c * (cw + gap); y = top + r * (rh + IN(0.26))
        box = rect(slide, x, y, cw, rh, fill=WHITE, line=LINE, lw=1)
        rect(slide, x, y, cw, Pt(3), fill=EMER)            # 上緣強調線
        text(slide, x + IN(0.26), y + IN(0.2), cw - IN(0.5), IN(0.5), emoji, size=22)
        text(slide, x + IN(0.26), y + IN(0.66), cw - IN(0.5), IN(0.4), title,
             size=15.5, color=CHAR, font=SANS, bold=True)
        text(slide, x + IN(0.26), y + IN(1.04), cw - IN(0.5), rh - IN(1.1), bodytx,
             size=11.5, color=BODY, font=SANS, sp=1.12)


def stacklist(slide, items, x, y, w, ih=IN(0.92)):
    for i, (ix, tt, sub) in enumerate(items):
        yy = y + i * ih
        text(slide, x, yy + IN(0.04), IN(1.0), IN(0.5), ix,
             size=15, color=EMER, font=SERIF, bold=True)
        text(slide, x + IN(1.0), yy, w - IN(1.0), IN(0.4), tt,
             size=14.5, color=CHAR, font=SANS, bold=True)
        text(slide, x + IN(1.0), yy + IN(0.38), w - IN(1.0), IN(0.4), sub,
             size=12, color=BODY, font=SANS)
        rect(slide, x, yy + ih - IN(0.12), w, Pt(0.75), fill=LINE)


def note(slide, x, y, w, h, tag, body, accent=EMER):
    rect(slide, x, y, w, h, fill=PAPER, line=LINE, lw=1)
    rect(slide, x, y, Pt(4), h, fill=accent)
    rect(slide, x + IN(0.25), y + IN(0.2), IN(1.5), IN(0.34), fill=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(slide, x + IN(0.25), y + IN(0.22), IN(1.5), IN(0.3), tag,
         size=10.5, color=WHITE, font=SANS, bold=True, align=PP_ALIGN.CENTER)
    text(slide, x + IN(0.28), y + IN(0.66), w - IN(0.56), h - IN(0.8), body,
         size=12.5, color=CHAR, font=SANS, sp=1.16)


# ===================== 投影片 =====================
def build():
    prs = Presentation(); prs.slide_width = EMW; prs.slide_height = EMH

    # 1 COVER
    s = blank(prs, WARM)
    rect(s, 0, 0, EMW, EMH, fill=EMER)
    rect(s, 0, 0, EMW, EMH, fill=RGBColor.from_string("053D2E"))  # 漸層感（疊色）
    frame(s, RGBColor.from_string("18654E"))
    rect(s, EMW - IN(3.4), IN(0.66), IN(2.5), IN(0.46), line=RGBColor.from_string("3E8C72"),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, EMW - IN(3.4), IN(0.7), IN(2.5), IN(0.4), "說明會簡報 · 範本",
         size=11.5, color=RGBColor.from_string("A7F3D0"), bold=True, align=PP_ALIGN.CENTER)
    text(s, IN(0.95), IN(2.2), IN(10), IN(0.4), "都市更新 · 危老重建　住戶說明會",
         size=13, color=RGBColor.from_string("A7F3D0"), bold=True, ls=2.4)
    text(s, IN(0.92), IN(2.7), IN(11.5), IN(2.0),
         "與您一起，把老屋變成安心的家", size=46, color=WHITE, font=SERIF, bold=True, sp=1.05)
    text(s, IN(0.95), IN(4.5), IN(9.6), IN(1.0),
         "本簡報為制式架構範本，帶您一次看懂：為何更新、如何進行、權益如何保障，以及您每一步該確認什麼。",
         size=15, color=MIST, sp=1.3)
    meta = [("案　名", "〔案名／門牌範圍〕"), ("實施者", "〔更新會／實施者〕"),
            ("日　期", "〔YYYY/MM/DD〕"), ("場　次", "〔第 N 次說明會〕")]
    for i, (k, v) in enumerate(meta):
        x = IN(0.95) + i * IN(2.95)
        text(s, x, IN(5.9), IN(2.8), IN(0.3), k, size=11, color=RGBColor.from_string("9FBAAC"))
        text(s, x, IN(6.2), IN(2.8), IN(0.4), v, size=14, color=WHITE, bold=True)

    # 2 AGENDA
    s = blank(prs); frame(s, LINE); header(s, "Agenda · 本日議程", "今天，我們會談這些")
    left = [("01", "為什麼要更新", "老屋風險與更新效益"),
            ("02", "基地與權屬現況", "區位、面積、權屬概況"),
            ("03", "法規路徑與容積", "都更／危老／防災、同意門檻、獎勵"),
            ("04", "規劃願景", "量體、戶型、公設與停車")]
    right = [("05", "權利變換與分回", "共同負擔、分回比、常見迷思"),
             ("06", "時程藍圖", "從評估到交屋的里程碑"),
             ("07", "您的保障", "信託、續建機制、合約與透明"),
             ("08", "常見問答與下一步", "Q&A、同意書、聯絡窗口")]
    stacklist(s, left, IN(0.9), IN(2.9), IN(5.6))
    stacklist(s, right, IN(6.9), IN(2.9), IN(5.6))
    footer(s, "都市更新／危老重建 · 住戶說明會")

    # 3 DIVIDER 01
    divider(prs, "01", "Why Renewal", "為什麼要更新",
            "更新不只是換新房子，而是把「安全、機能、資產價值」一次補齊。")

    # 4 WHY
    s = blank(prs); frame(s, LINE); header(s, "老屋的隱形成本", "不更新，風險會一直累積")
    cards(s, [
        ("🏚️", "耐震不足", "早期建築標準較低，地震時結構風險高，補強有極限。"),
        ("🚿", "管線老化", "漏水、壁癌、消防與電力負載不符現代需求。"),
        ("🚗", "機能落後", "無電梯、無車位、無無障礙，居住品質與轉手性下降。"),
        ("📉", "資產停滯", "屋齡越高、增值越慢；更新可重塑資產價值。"),
        ("🛡️", "安全升級", "新制耐震、消防、無障礙、智慧管理一次到位。"),
        ("🌱", "世代資產", "把老屋轉為可傳承、好出租、好自住的新資產。"),
    ])
    footer(s, "為什麼要更新")

    # 5 SITE
    s = blank(prs); frame(s, LINE); header(s, "基地現況 · Site", "我們的基地，條件如何？")
    stacklist(s, [
        ("區位", "〔行政區／路段〕", "生活機能、交通與學區概述"),
        ("面積", "〔基地面積 ㎡／坪〕", "更新單元範圍與地號筆數"),
        ("分區", "〔使用分區／容積率〕", "可建範圍的法定基礎"),
        ("權屬", "〔所有權人戶數〕", "同意整合的分母"),
    ], IN(0.9), IN(2.9), IN(6.2))
    note(s, IN(7.5), IN(2.9), IN(4.9), IN(2.6), "透明原則",
         "本頁所有數字以地政與使照文件為準，現場提供查詢方式。\n\n我們承諾：對外說明、合約附件、送審文件——同一套數字、同一個版本。")
    footer(s, "基地與權屬現況")

    # 6 DIVIDER 02
    divider(prs, "02", "Law & Floor Area", "法規路徑與容積",
            "走哪條路、要多少人同意、能爭取多少獎勵——先把規則講清楚。")

    # 7 PATH
    s = blank(prs); frame(s, LINE); header(s, "三條路徑，怎麼選", "都更 · 危老 · 防災")
    cards(s, [
        ("都市更新", "多數決整合", "需達法定同意比例（依劃定方式而定），審議程序較完整，容積獎勵空間較大。"),
        ("危老重建", "全體同意", "需 100% 同意、程序較快，適合戶數單純的基地；無須劃定更新單元。"),
        ("防災型", "時程獎勵", "具時程獎勵但有申請落日期限，越早送件越有利。"),
    ], rh=IN(2.0))
    note(s, IN(0.9), IN(5.3), IN(11.5), IN(1.1), "提醒",
         "同意門檻與容積獎勵以現行條文為準，本頁為路徑概念說明；個案適用條件由專業團隊逐項查證後提供書面對照。")
    footer(s, "法規路徑與容積")

    # 8 DIVIDER 03
    divider(prs, "03", "Design Vision", "規劃願景",
            "把基地條件，轉化成安全、好住、保值的新家。")

    # 9 DESIGN
    s = blank(prs); frame(s, LINE); header(s, "規劃重點 · Planning", "新家會長這樣")
    stacklist(s, [
        ("量體", "〔地上 N 層／地下 N 層〕", "退縮、棟距與日照通風"),
        ("戶型", "〔坪型配置〕", "分回與銷售戶型規劃"),
        ("公設", "〔公設比／公設項目〕", "大廳、電梯、無障礙、社區設施"),
        ("停車", "〔車位數／型式〕", "平面／機械、車道規劃"),
    ], IN(0.9), IN(2.9), IN(6.2))
    note(s, IN(7.5), IN(2.9), IN(4.9), IN(2.6), "意象示意",
         "此處置入外觀透視、平面配置、公設意象圖。\n\n提醒：示意圖以核定圖說為準，最終以建照與權變核定內容為依據。")
    footer(s, "規劃願景")

    # 10 DIVIDER 04
    divider(prs, "04", "Rights & Allocation", "權利變換與分回",
            "最關鍵的一頁：您出地、團隊出力，蓋好後怎麼分。")

    # 11 ALLOCATION（比例帶，原生圖形）
    s = blank(prs); frame(s, LINE); header(s, "分回的邏輯 · 示意", "總價值，怎麼分？")
    STONE_BAR = RGBColor.from_string("E0DCD2")
    bx, by, bw, bh = IN(0.9), IN(2.95), IN(11.5), IN(1.0)
    rect(s, bx, by, bw, bh, fill=STONE_BAR, line=LINE, lw=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, bx, by, Emu(int(bw * 0.6)), bh, fill=EMER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, bx, by, Emu(int(bw * 0.6)), bh, "共同負擔　60%",
         size=16, color=WHITE, font=SANS, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, bx + Emu(int(bw * 0.6)), by, Emu(int(bw * 0.4)), bh, "住戶分回　40%",
         size=16, color=CHAR, font=SANS, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, bx, by + IN(1.15), bw, IN(0.4), "示意 · 非個案數字（實際比例依個案共同負擔而定）",
         size=12.5, color=EMER, font=SERIF, bold=True, align=PP_ALIGN.CENTER)
    legend = [(EMER, "共同負擔", "工程・管維・權變・利息・稅捐・管理（六大科目）"),
              (STONE_BAR, "住戶分回", "＝ 1 − 共同負擔比；經估價師查核、主管機關核定")]
    for i, (col, tt, sub) in enumerate(legend):
        x = IN(0.9) + i * IN(5.9)
        rect(s, x, IN(4.55), IN(0.3), IN(0.3), fill=col, line=LINE)
        text(s, x + IN(0.45), IN(4.5), IN(5.2), IN(0.4), tt, size=15.5, color=CHAR, font=SERIF, bold=True)
        text(s, x + IN(0.45), IN(4.9), IN(5.2), IN(0.6), sub, size=12, color=BODY, sp=1.1)
    note(s, IN(0.9), IN(5.7), IN(11.5), IN(1.0), "破解迷思",
         "「一坪換一坪」是迷思——分回是「總價值 × 分回比」的結果，不是固定坪數承諾。實際分回以權利變換核定為準。")
    footer(s, "權利變換與分回")

    # 12 TIMELINE
    s = blank(prs); frame(s, LINE); header(s, "時程藍圖 · Roadmap", "從評估到交屋，分四大階段")
    phases = [
        (EMER,                          "前期評估 S1–S4", ["S1 基地權屬調查", "S2 法規容積試算", "S3 坪效量體規劃", "S4 投報財務評估"]),
        (EMER_L,                        "整合 S5–S6",    ["S5 說明會與意願整合", "S6 簽約"]),
        (RGBColor.from_string("0F766E"),"法定程序 S7–S8", ["S7 事業／權變送審", "S8 審議・公展・核定"]),
        (RGBColor.from_string("047857"),"執行 S9–S11",   ["S9 設計與發包", "S10 施工", "S11 銷售・交屋・管委會"]),
    ]
    gap = IN(0.2); cw = (EMW - IN(1.8) - gap * 3) / 4
    for i, (col, lbl, steps) in enumerate(phases):
        x = IN(0.9) + i * (cw + gap)
        rect(s, x, IN(2.7), cw, IN(0.5), fill=col)
        text(s, x, IN(2.8), cw, IN(0.32), lbl, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        bh = IN(0.42) * len(steps) + IN(0.3)
        rect(s, x, IN(3.2), cw, bh, fill=WHITE, line=LINE, lw=1)
        for j, st in enumerate(steps):
            text(s, x + IN(0.18), IN(3.34) + j * IN(0.42), cw - IN(0.32), IN(0.4),
                 st, size=11.5, color=BODY)
    note(s, IN(0.9), IN(5.6), IN(11.5), IN(0.95), "我們在這裡",
         "目前階段：〔SN〕。每個階段都有明確產出物與檢核點，進度公開可查。")
    footer(s, "時程藍圖")

    # 13 DIVIDER 05
    divider(prs, "05", "Your Safeguards", "您的保障",
            "把信任建立在「制度」上，而不是任何一家公司的口頭承諾。")

    # 14 SAFEGUARDS
    s = blank(prs); frame(s, LINE); header(s, "四道防線 · Safeguards", "您的權益，靠制度保障")
    cards(s, [
        ("🏦", "信託專戶", "資金進信託、撥付綁里程碑，專款專用、不被挪用。"),
        ("🔁", "續建機制", "萬一實施者出狀況，信託／融資機構可啟動續建、接管工程，您的權益優先受保障。"),
        ("📜", "合約保障", "權利義務、違約責任、撤回與解約條件，白紙黑字、對等明確。"),
        ("🧾", "單一數源", "說明、合約、送審同版同源、日期戳記，逐戶分配可回算。"),
        ("⚖️", "獨立估價", "權變交由估價師查核、主管機關核定，非實施者自說自話。"),
        ("👀", "程序公開", "審議、公展、核定皆公開可查，重大事項留痕。"),
    ])
    footer(s, "您的保障")

    # 15 Q&A
    s = blank(prs); frame(s, LINE); header(s, "常見問答 · FAQ", "大家最常問的事")
    qa = [
        ("Q. 現在簽的是什麼？會不會被綁死？", "同意整合與正式合約是不同階段；授權範圍、撤回與解約條件都會逐條說明，您看懂再簽。"),
        ("Q. 萬一建商蓋到一半跑了？", "這正是信託＋續建機制要解決的問題：由制度接手續建，權益優先保障。"),
        ("Q. 真的能「一坪換一坪」嗎？", "分回是總價值×分回比的結果，視共同負擔而定；我們提供透明試算，不做不實承諾。"),
        ("Q. 我可以不參加嗎？", "更新採法定同意門檻；達門檻後依權利變換處理，過程有估價與安置等法定保障。"),
        ("Q. 數字會不會說一套做一套？", "我們承諾單一數源：對外、合約、送審同一版本，逐戶分配回算一致。"),
        ("Q. 要多久？現在到哪了？", "依四大階段推進，目前在〔SN〕；每階段產出與檢核點公開可查。"),
    ]
    cw = (EMW - IN(1.8) - IN(0.5)) / 2
    for i, (q, a) in enumerate(qa):
        r, c = divmod(i, 2)
        x = IN(0.9) + c * (cw + IN(0.5)); y = IN(2.75) + r * IN(1.35)
        rect(s, x, y + IN(0.04), Pt(3), IN(1.0), fill=EMER)
        text(s, x + IN(0.22), y, cw - IN(0.3), IN(0.5), q, size=14, color=EMER, font=SANS, bold=True)
        text(s, x + IN(0.22), y + IN(0.44), cw - IN(0.3), IN(0.8), a, size=12, color=BODY, sp=1.12)
    footer(s, "常見問答")

    # 16 NEXT STEPS
    s = blank(prs); frame(s, LINE); header(s, "下一步 · Next", "接下來，這樣進行")
    stacklist(s, [
        ("01", "領取資料包", "說明簡報、Q&A、同意書與條文對照表"),
        ("02", "個別諮詢", "一對一說明您的權屬與分回試算"),
        ("03", "看懂再簽", "確認授權範圍與保障條款後，再決定"),
        ("04", "持續更新", "進度、審議結果、財務調整即時公開"),
    ], IN(0.9), IN(2.9), IN(6.2))
    note(s, IN(7.5), IN(2.9), IN(4.9), IN(2.6), "給您的承諾",
         "我們寧可多花時間把疑慮講清楚，也不催促您簽下看不懂的文件。\n\n專業的質疑，是把案子做好的養分。", accent=EMER_L)
    footer(s, "下一步")

    # 17 CLOSING
    s = blank(prs); rect(s, 0, 0, EMW, EMH, fill=RGBColor.from_string("053D2E"))
    frame(s, RGBColor.from_string("18654E"))
    text(s, IN(1), IN(2.5), IN(11.33), IN(0.4), "THANK YOU", size=13,
         color=RGBColor.from_string("A7F3D0"), bold=True, ls=3.0, align=PP_ALIGN.CENTER)
    text(s, IN(1), IN(3.0), IN(11.33), IN(1.0), "謝謝您的參與", size=44, color=WHITE,
         font=SERIF, bold=True, align=PP_ALIGN.CENTER)
    rect(s, EMW/2 - IN(0.35), IN(4.2), IN(0.7), Pt(3), fill=EMER_L)
    text(s, IN(1.8), IN(4.5), IN(9.73), IN(0.8),
         "您的每一個問題，都讓這個案子更穩健。歡迎於會後一對一諮詢。",
         size=15, color=MIST, align=PP_ALIGN.CENTER, sp=1.3)
    contact = [("窗口", "〔整合專員／電話〕"), ("實施者", "〔更新會／實施者〕"), ("諮詢時間", "〔時段／地點〕")]
    for i, (k, v) in enumerate(contact):
        x = IN(2.4) + i * IN(2.95)
        text(s, x, IN(5.7), IN(2.8), IN(0.3), k, size=11, color=RGBColor.from_string("9FBAAC"), align=PP_ALIGN.CENTER)
        text(s, x, IN(6.0), IN(2.8), IN(0.4), v, size=13, color=RGBColor.from_string("A7F3D0"), bold=True, align=PP_ALIGN.CENTER)

    out = "都更說明會簡報範本.pptx"
    prs.save(out)
    print(f"✓ 已輸出 {out}（{len(prs.slides.__iter__.__self__._sldIdLst)} 張投影片）")
    return out


def divider(prs, num, kicker, title, sub):
    s = blank(prs); rect(s, 0, 0, EMW, EMH, fill=EMER)
    rect(s, 0, 0, EMW, EMH, fill=RGBColor.from_string("06513D"))
    frame(s, RGBColor.from_string("18654E"))
    text(s, IN(0.95), IN(1.5), IN(4), IN(1.6), num, size=92, color=RGBColor.from_string("0C6149"),
         font=SERIF, bold=True)
    text(s, IN(1.0), IN(3.3), IN(10), IN(0.4), kicker.upper(), size=12.5,
         color=RGBColor.from_string("A7F3D0"), bold=True, ls=2.8)
    text(s, IN(0.95), IN(3.8), IN(11), IN(1.2), title, size=46, color=WHITE, font=SERIF, bold=True)
    text(s, IN(0.98), IN(5.2), IN(9.5), IN(1.0), sub, size=15.5, color=MIST, sp=1.35)


if __name__ == "__main__":
    build()
